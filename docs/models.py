from django.db import models
from django.shortcuts import reverse
from ckeditor.fields import RichTextField
from django_ckeditor_5.fields import CKEditor5Field

# дополнения для pptx и pdf
import os
import logging
from django.conf import settings
from django.core.files.storage import default_storage
from django.core.cache import cache
from django.db.models.signals import post_save
from django.dispatch import receiver
from PIL import Image
import io
import time
from pathlib import Path


class Direction_of_business(models.Model):
    # направление бизнеса (глобальные)
    name = models.CharField(max_length=100, verbose_name="Название бизнеса")
    description = models.TextField(blank=True, verbose_name="Описание")

    def __str__(self):
        return self.name

    class Meta:
        verbose_name = "Направления бизнеса"
        verbose_name_plural = "Направления бизнеса"


class Direction(models.Model):
    # направления (локальные)
    name = models.CharField(max_length=100, verbose_name="Название рубрики")
    description = models.TextField(blank=True, verbose_name="Описание")
    global_section = models.ForeignKey(Direction_of_business, on_delete=models.CASCADE,
                                       verbose_name="Направление бизнеса")

    def __str__(self):
        return self.name

    class Meta:
        verbose_name = "Направление"
        verbose_name_plural = "Направления"


class Section(models.Model):
    # рубрики
    name = models.CharField(max_length=100, verbose_name="Название рубрики")
    description = models.TextField(blank=True, verbose_name="Описание")
    global_section = models.ForeignKey(Direction_of_business, on_delete=models.CASCADE,
                                       verbose_name="Направление бизнеса")

    def __str__(self):
        return self.name

    class Meta:
        verbose_name = "Рубрика"
        verbose_name_plural = "Рубрики"


class Doc_class(models.Model):
    # типы документов
    name = models.CharField(max_length=100, verbose_name="Название категории")
    description = models.TextField(blank=True, verbose_name="Описание")
    # category = models.ForeignKey(Section, on_delete=models.CASCADE, verbose_name="Рубрика")

    def __str__(self):
        return self.name

    class Meta:
        verbose_name = "Тип документа"
        verbose_name_plural = "Типы документов"


# дополнения для pptx и pdf
logger = logging.getLogger(__name__)


def document_upload_to(instance, filename):
    return f'documents/{int(time.time())}_{filename}'


class Document(models.Model):
    # документ
    STATUS_CHOICES = [
        ('active', 'Действующая редакция'),
        ('obsolete', 'Утратил силу'),
    ]

    title = models.CharField(max_length=255, verbose_name="Название документа")
    file = models.FileField(upload_to='documents/', verbose_name="Файл документа", blank=True, null=True, )
    status = models.CharField(max_length=10, choices=STATUS_CHOICES, default='active', verbose_name="Статус")
    created_at = models.DateTimeField(auto_now_add=True, verbose_name="Дата создания")
    updated_at = models.DateTimeField(auto_now=True, verbose_name="Дата обновления")
    global_section = models.ForeignKey(Direction_of_business, on_delete=models.CASCADE, verbose_name="Направление бизнеса")
    section = models.ForeignKey(Direction, on_delete=models.CASCADE, blank=True, null=True, verbose_name="Направление", default='none')
    category = models.ForeignKey(Section, on_delete=models.CASCADE, blank=True, null=True, verbose_name="Рубрика")
    subcategory = models.ForeignKey(Doc_class, on_delete=models.CASCADE, verbose_name="Тип документа")
    is_template = models.BooleanField(default=False, verbose_name="Это шаблон")
    description = CKEditor5Field(verbose_name='Полное описание', config_name='extends', blank=True, null=True,)

    def __str__(self):
        return self.title

    class Meta:
        verbose_name = "Документ"
        verbose_name_plural = "Документы"

    def get_absolute_url(self):
        return reverse('knowledge_base:document_detail', args=[self.id])

    # дополнения для pptx и pdf
    def save(self, *args, **kwargs):
        # Определяем тип файла перед сохранением
        if self.file.name.lower().endswith('.pptx'):
            self.file_type = 'pptx'
        elif self.file.name.lower().endswith('.pdf'):
            self.file_type = 'pdf'
        else:
            self.file_type = 'other'

        super().save(*args, **kwargs)

    def get_cache_key(self):
        return f'document_{self.id}_slides'

    def convert_file(self):
        try:
            if self.file_type == 'pptx':
                return self._convert_pptx()
            elif self.file_type == 'pdf':
                return self._convert_pdf()
            return None
        except Exception as e:
            logger.error(f"Error converting file {self.file.name}: {str(e)}")
            self.conversion_error = str(e)
            self.save()
            return None

    def _convert_pptx(self):
        try:
            from pptx import Presentation

            pptx_path = self.file.path
            pres = Presentation(pptx_path)
            base_dir = Path(settings.MEDIA_ROOT) / 'converted' / str(self.id)
            base_dir.mkdir(parents=True, exist_ok=True)

            slide_images = []

            for i, slide in enumerate(pres.slides):
                img_path = base_dir / f'slide_{i + 1}.jpg'

                # В реальном приложении здесь должна быть конвертация слайда в изображение
                # Для примера создаем пустое изображение
                img = Image.new('RGB', (800, 600), color='white')
                img.save(img_path)

                slide_data = {
                    'number': i + 1,
                    'image_url': f'/media/converted/{self.id}/slide_{i + 1}.jpg',
                    'notes': slide.notes_slide.notes_text_frame.text if slide.notes_slide else ''
                }
                slide_images.append(slide_data)

            # Кэшируем результат
            cache_key = self.get_cache_key()
            cache.set(cache_key, slide_images, timeout=60 * 60 * 24)  # Кэш на 24 часа

            self.converted = True
            self.save()
            return slide_images

        except Exception as e:
            logger.error(f"PPTX conversion error: {str(e)}")
            raise

    def _convert_pdf(self):
        try:
            from pdf2image import convert_from_path

            pdf_path = self.file.path
            base_dir = Path(settings.MEDIA_ROOT) / 'converted' / str(self.id)
            base_dir.mkdir(parents=True, exist_ok=True)

            images = convert_from_path(pdf_path, dpi=200, fmt='jpeg')
            slide_images = []

            for i, image in enumerate(images):
                img_path = base_dir / f'slide_{i + 1}.jpg'
                image.save(img_path, 'JPEG')

                slide_data = {
                    'number': i + 1,
                    'image_url': f'/media/converted/{self.id}/slide_{i + 1}.jpg',
                    'notes': f'Страница {i + 1}'
                }
                slide_images.append(slide_data)

            # Кэшируем результат
            cache_key = self.get_cache_key()
            cache.set(cache_key, slide_images, timeout=60 * 60 * 24)  # Кэш на 24 часа

            self.converted = True
            self.save()
            return slide_images

        except Exception as e:
            logger.error(f"PDF conversion error: {str(e)}")
            raise


@receiver(post_save, sender=Document)
def process_document(sender, instance, created, **kwargs):
    if created and instance.file_type in ['pptx', 'pdf']:
        try:
            instance.convert_file()
        except Exception as e:
            logger.error(f"Error processing document: {str(e)}")